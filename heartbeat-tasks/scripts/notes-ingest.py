#!/usr/bin/env python3
"""Ingest GoodNotes PDF exports into vault as structured course notes.

For each new/unprocessed PDF under Cal Poly Pomona:
1. Determine course from folder path (deterministic)
2. Call Claude Sonnet to read the PDF and extract key concepts
3. Write structured markdown to vault/shared/course-notes/<course>/
4. Track ingestion state to avoid reprocessing

Runs as a heartbeat script every 4 hours.
"""

import os
import sys
import json
import subprocess
import datetime
import re

HARNESS_ROOT = os.environ.get(
    "HARNESS_ROOT",
    os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
)

# LLM provider — defaults to claude-cli, overridable via LLM_PROVIDER env var
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from lib.llm_provider import get_provider, get_default_model, LLMError
TASKS_DIR = os.path.join(HARNESS_ROOT, "heartbeat-tasks")
STATE_FILE = os.path.join(TASKS_DIR, "notes-ingest.state.json")
KNOWN_FILES_STATE = os.path.join(TASKS_DIR, "goodnotes-watch.state-files.json")
VAULT_DIR = os.path.join(HARNESS_ROOT, "vault", "shared", "course-notes")
NOTIFY_FILE = os.path.join(TASKS_DIR, "pending-notifications.jsonl")

# Max PDFs to process per run (cost control)
MAX_PER_RUN = 3

# Map GoodNotes folder names → vault directory names and Discord channels
# Loaded from course-map.json (gitignored) — copy course-map.example.json to get started
_course_map_path = os.path.join(TASKS_DIR, "course-map.json")
if os.path.exists(_course_map_path):
    with open(_course_map_path) as _f:
        COURSE_MAP = json.load(_f).get("goodnotes", {})
else:
    COURSE_MAP = {}


def _find_google_drive():
    """Auto-detect the Google Drive CloudStorage mount."""
    cloud_dir = os.path.expanduser("~/Library/CloudStorage")
    if os.path.isdir(cloud_dir):
        for entry in os.listdir(cloud_dir):
            if entry.startswith("GoogleDrive-"):
                return os.path.join(cloud_dir, entry, "My Drive", "GoodNotes")
    return None


def load_state():
    if os.path.exists(STATE_FILE):
        with open(STATE_FILE) as f:
            state = json.load(f)
        # Migrate: add failures tracking if missing
        if "failures" not in state:
            state["failures"] = {}
        return state
    return {"ingested": [], "failures": {}, "last_run": None}

MAX_FAILURES = 3  # Skip PDFs that fail this many times


def save_state(state):
    state["last_run"] = datetime.datetime.now().isoformat()
    tmp = STATE_FILE + ".tmp"
    with open(tmp, "w") as f:
        json.dump(state, f, indent=2)
    os.rename(tmp, STATE_FILE)


def notify(message, channel="goodnotes"):
    entry = {
        "task": "notes-ingest",
        "channel": channel,
        "summary": message,
        "timestamp": datetime.datetime.now().isoformat(),
    }
    with open(NOTIFY_FILE, "a") as f:
        f.write(json.dumps(entry) + "\n")


def parse_course(rel_path):
    """Extract course info from GoodNotes relative path.

    Paths look like: Cal Poly Pomona/Numerical Methods/Notes/2-24.pdf
    Returns course map entry or None.
    """
    parts = rel_path.split("/")
    if len(parts) < 2:
        return None
    if "Cal Poly" not in parts[0]:
        return None
    folder_name = parts[1]
    return COURSE_MAP.get(folder_name)


def make_vault_filename(rel_path, course_info):
    """Generate a vault-friendly filename from the PDF path."""
    basename = os.path.splitext(os.path.basename(rel_path))[0]
    # Determine subfolder context (Notes, Homework, Discussions, etc.)
    parts = rel_path.split("/")
    subfolder = ""
    if len(parts) >= 3:
        # e.g., "Cal Poly Pomona/Numerical Methods/Notes/2-24.pdf" → subfolder = "notes"
        potential = parts[2].lower()
        if potential in ("notes", "homework", "discussions", "quizzes"):
            subfolder = potential + "-"

    # Clean up the basename for a filename
    clean = re.sub(r'[^\w\s-]', '', basename).strip()
    clean = re.sub(r'\s+', '-', clean).lower()
    return f"{subfolder}{clean}"


def extract_pptx_text(pptx_path):
    """Extract text from a PowerPoint file using python-pptx.
    Returns (slides_text, has_substantial_text) tuple.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("  python-pptx not installed — skipping PPTX extraction", file=sys.stderr)
        return None, False

    prs = Presentation(pptx_path)
    slides_text = []
    total_chars = 0

    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## Slide {i}"]

        # Extract title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_lines.append(f"**{slide.shapes.title.text.strip()}**\n")

        # Extract all text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_lines.append(text)
                        total_chars += len(text)

            # Extract table data
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    slide_lines.append("| " + " | ".join(cells) + " |")
                    total_chars += sum(len(c) for c in cells)

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"\n*Speaker notes:* {notes}")
                total_chars += len(notes)

        slides_text.append("\n".join(slide_lines))

    full_text = "\n\n".join(slides_text)
    # Consider "substantial" if average >30 chars per slide
    has_substantial_text = total_chars > (len(prs.slides) * 30) if prs.slides else False

    return full_text, has_substantial_text


def ingest_pptx(pptx_path, rel_path, course_info):
    """Process a PowerPoint file — extract text directly or use Claude for visual slides."""
    vault_dir = os.path.join(VAULT_DIR, course_info["vault_dir"])
    os.makedirs(vault_dir, exist_ok=True)

    vault_filename = make_vault_filename(rel_path, course_info)
    vault_path = os.path.join(vault_dir, f"{vault_filename}.md")

    if os.path.exists(vault_path):
        return vault_path

    # Determine content type
    parts = rel_path.split("/")
    content_type = "lecture slides"
    if len(parts) >= 3:
        sub = parts[2].lower()
        if "homework" in sub or "hw" in sub:
            content_type = "homework"
        elif "discussion" in sub:
            content_type = "discussion"

    # Try text extraction first
    extracted_text, has_text = extract_pptx_text(pptx_path)

    if has_text and extracted_text:
        # Text-heavy slides — use Claude to structure the extracted text
        prompt = f"""Here is text extracted from PowerPoint slides for a {course_info['display']} class ({content_type}).
Structure this into clean, organized study notes.

Extracted slide content:
{extracted_text[:8000]}

Output as clean markdown with:
1. **Topic/Title** — what these slides cover
2. **Key Concepts** — main ideas, definitions, theorems
3. **Formulas/Equations** — any mathematical formulas (use LaTeX notation)
4. **Examples** — worked examples or practice problems
5. **Key Takeaways** — what a student should remember

Be thorough but concise. Do NOT include slide numbers or "Slide X" headers — reorganize by topic."""

        try:
            llm = get_provider()
            response = llm.complete(prompt, model=get_default_model(), timeout=120, max_turns=3)
            content = response.text

            if content and len(content) > 50:
                    frontmatter = f"""---
course: {course_info['display']}
source: {rel_path}
type: {content_type}
format: pptx
ingested: {datetime.datetime.now().strftime('%Y-%m-%d')}
---

"""
                    with open(vault_path, "w") as f:
                        f.write(frontmatter + content)
                    return vault_path
        except Exception as e:
            print(f"  LLM structuring failed: {e}", file=sys.stderr)

        # Fallback: write raw extracted text
        frontmatter = f"""---
course: {course_info['display']}
source: {rel_path}
type: {content_type}
format: pptx-raw
ingested: {datetime.datetime.now().strftime('%Y-%m-%d')}
---

"""
        with open(vault_path, "w") as f:
            f.write(frontmatter + extracted_text)
        return vault_path

    else:
        # Visual/formula-heavy slides — use Claude with Read tool (same as PDF path)
        return ingest_pdf(pptx_path, rel_path, course_info)


def ingest_pdf(pdf_path, rel_path, course_info):
    """Call Claude Sonnet to read a PDF and extract structured notes."""
    vault_dir = os.path.join(VAULT_DIR, course_info["vault_dir"])
    os.makedirs(vault_dir, exist_ok=True)

    vault_filename = make_vault_filename(rel_path, course_info)
    vault_path = os.path.join(vault_dir, f"{vault_filename}.md")

    # Skip if already exists in vault
    if os.path.exists(vault_path):
        return vault_path

    # Determine content type from subfolder
    parts = rel_path.split("/")
    content_type = "lecture notes"
    if len(parts) >= 3:
        sub = parts[2].lower()
        if "homework" in sub or "hw" in sub:
            content_type = "homework"
        elif "discussion" in sub:
            content_type = "discussion"
        elif "quiz" in sub:
            content_type = "quiz preparation"

    prompt = f"""Use the Read tool to open this PDF file, then extract structured notes from it.

IMPORTANT: You MUST use the Read tool with this exact file path to read the PDF:
{pdf_path}

If the PDF is large (many pages), read it in chunks using the pages parameter (e.g., pages="1-10", then pages="11-20"). Maximum 10 pages per Read call to avoid memory issues.

This is {content_type} from a {course_info['display']} class. After reading the PDF, extract:
1. **Topic/Title** — what this covers
2. **Key Concepts** — main ideas, definitions, theorems
3. **Formulas/Equations** — any mathematical formulas (use LaTeX notation)
4. **Examples** — worked examples or practice problems
5. **Key Takeaways** — what a student should remember

Output as clean markdown. Be thorough but concise. If handwriting is unclear, note it as [UNCLEAR].
Do NOT include any meta-commentary — just the extracted content."""

    try:
        llm = get_provider()
        response = llm.complete(
            prompt, model=get_default_model(), timeout=600,
            max_turns=15, allowed_tools=["Read"],
        )
        content = response.text

        if not content or len(content) < 50:
            print(f"  Empty or too short response for {rel_path}", file=sys.stderr)
            return None

        # Write vault entry with frontmatter
        frontmatter = f"""---
course: {course_info['display']}
source: {rel_path}
type: {content_type}
ingested: {datetime.datetime.now().strftime('%Y-%m-%d')}
---

"""
        with open(vault_path, "w") as f:
            f.write(frontmatter + content)

        return vault_path

    except LLMError as e:
        print(f"  LLM error: {e}", file=sys.stderr)
        return None
    except Exception as e:
        print(f"  Error: {e}", file=sys.stderr)
        return None


def main():
    drive_dir = _find_google_drive()
    if not drive_dir:
        print("Google Drive not found")
        sys.exit(1)

    # Load state
    state = load_state()
    ingested = set(state.get("ingested", []))

    # Load known files from goodnotes-watch
    if not os.path.exists(KNOWN_FILES_STATE):
        print("No goodnotes-watch state file — run goodnotes-watch first")
        sys.exit(0)

    with open(KNOWN_FILES_STATE) as f:
        known_files = json.load(f)

    # Filter to Cal Poly Pomona academic files only (PDF + PPTX)
    SUPPORTED_EXTENSIONS = (".pdf", ".pptx")
    failures = state.get("failures", {})
    to_process = []
    for rel_path in sorted(known_files):
        if rel_path in ingested:
            continue
        # Only process supported file types
        if not rel_path.lower().endswith(SUPPORTED_EXTENSIONS):
            continue
        # Skip files that have failed too many times
        if failures.get(rel_path, 0) >= MAX_FAILURES:
            continue
        course_info = parse_course(rel_path)
        if not course_info:
            continue
        full_path = os.path.join(drive_dir, rel_path)
        if not os.path.exists(full_path):
            continue
        to_process.append((full_path, rel_path, course_info))

    if not to_process:
        print("All academic PDFs already ingested")
        save_state(state)
        return

    # Cap per run
    batch = to_process[:MAX_PER_RUN]
    print(f"Processing {len(batch)} of {len(to_process)} pending files...")

    results = {"success": [], "failed": []}
    for full_path, rel_path, course_info in batch:
        print(f"  Ingesting: {rel_path}")
        # Route to appropriate handler based on file type
        if rel_path.lower().endswith(".pptx"):
            vault_path = ingest_pptx(full_path, rel_path, course_info)
        else:
            vault_path = ingest_pdf(full_path, rel_path, course_info)
        if vault_path:
            ingested.add(rel_path)
            results["success"].append((rel_path, course_info))
            print(f"    → {os.path.basename(vault_path)}")
        else:
            failures[rel_path] = failures.get(rel_path, 0) + 1
            results["failed"].append(rel_path)
            remaining_tries = MAX_FAILURES - failures[rel_path]
            print(f"    → FAILED ({remaining_tries} retries left)")

    # Save state
    state["ingested"] = sorted(ingested)
    state["failures"] = failures
    save_state(state)

    # Notify per course
    if results["success"]:
        from collections import defaultdict
        by_course = defaultdict(list)
        for rel_path, course_info in results["success"]:
            by_course[course_info["channel"]].append(os.path.basename(rel_path))

        for channel, files in by_course.items():
            course_display = next(
                c["display"] for c in COURSE_MAP.values() if c["channel"] == channel
            )
            lines = [f"**{len(files)} note(s) ingested for {course_display}**\n"]
            for f in files:
                lines.append(f"  • {f}")
            lines.append(f"\nUse `/academics study <topic>` to generate study material.")
            notify("\n".join(lines), channel=channel)

    remaining = len(to_process) - len(batch)
    if remaining > 0:
        print(f"{remaining} PDFs remaining — will process on next run")

    print(f"Done: {len(results['success'])} ingested, {len(results['failed'])} failed")


if __name__ == "__main__":
    main()
